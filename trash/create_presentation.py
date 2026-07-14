from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Slide 1: Title
slide_layout = prs.slide_layouts[0] # Title Slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "HubBot"
subtitle.text = "Гибридная система делегирования задач\nTelegram + Веб-платформа"

# Slide 2: Проблема
slide_layout = prs.slide_layouts[1] # Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Проблема: Хаос в задачах"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "С чем сталкивается руководитель:"
p = tf.add_paragraph()
p.text = "Задачи теряются в переписках Telegram."
p = tf.add_paragraph()
p.text = "Нет контроля сроков и дедлайнов."
p = tf.add_paragraph()
p.text = "Отчеты сотрудников смешиваются с обычными сообщениями."

# Slide 3: Решение
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Решение: HubBot"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Разделение рабочих зон:"
p = tf.add_paragraph()
p.text = "Для руководителя: Умный Telegram-бот. Удобно ставить задачи с телефона."
p = tf.add_paragraph()
p.text = "Для сотрудников: Профессиональный веб-портал. Удобно видеть тикеты и загружать отчеты."
p = tf.add_paragraph()
p.text = "Единая база данных: Изменения на сайте мгновенно синхронизируются с ботом."

# Slide 4: Telegram-Бот
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Telegram-Бот (Панель управления)"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Возможности бота:"
p = tf.add_paragraph()
p.text = "Создание задач (/add) через многошаговый диалог."
p = tf.add_paragraph()
p.text = "Установка дедлайнов на принятие и на выполнение."
p = tf.add_paragraph()
p.text = "Просмотр всех свободных задач (/tasks)."
p = tf.add_paragraph()
p.text = "Получение файлов с результатами напрямую от веб-сервера."

# Slide 5: Веб-Интерфейс
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Веб-Интерфейс (Диспетчерская)"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Рабочее место сотрудника:"
p = tf.add_paragraph()
p.text = "Строгий дизайн в стиле «Диспетчерской тикетов»."
p = tf.add_paragraph()
p.text = "Доска задач с распределением статусов (Свободно, В процессе, Выполнено)."
p = tf.add_paragraph()
p.text = "Собственный профиль с историей выполненных заданий."
p = tf.add_paragraph()
p.text = "Загрузка файлов в 1 клик с красивым отображением имени."

# Slide 6: Автоматизация и Контроль
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Жесткий контроль сроков"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Система не прощает просрочек:"
p = tf.add_paragraph()
p.text = "Фоновый процесс (Cleaner Thread) работает 24/7."
p = tf.add_paragraph()
p.text = "Каждую минуту проверяются дедлайны по всем задачам."
p = tf.add_paragraph()
p.text = "Если задача просрочена — она автоматически удаляется."
p = tf.add_paragraph()
p.text = "Руководитель получает push-уведомление об удалении тикета."

# Slide 7: Технологический стек
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Технологический стек"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Что внутри HubBot?"
p = tf.add_paragraph()
p.text = "Бэкенд: Python (Flask) + Telebot."
p = tf.add_paragraph()
p.text = "База Данных: SQLite."
p = tf.add_paragraph()
p.text = "Фронтенд: HTML, JavaScript, CSS Grid."
p = tf.add_paragraph()
p.text = "Деплой: Docker (контейнеризация для быстрого развертывания)."

# Slide 8: Заключение
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Спасибо за внимание!"
subtitle.text = "HubBot — порядок в задачах, порядок в делах."

prs.save('HubBot_Presentation.pptx')
